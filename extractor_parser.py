import re
import sys
import glob
import os
from openpyxl import load_workbook
from pptx import Presentation


EMOJI_VOTE_MAP = {
    "\U0001f7e2": "YES",            # 🟢
    "\U0001f534": "NO",             # 🔴
    "\U0001f7e1": "NEEDS DISCUSSION",  # 🟡
}

BOILERPLATE_RE = re.compile(
    r"clear\s+yes|clear\s+no|needs?\s+discussion|strong\s+scores?|"
    r"strong\s+recommendations?|little\s+disagreement|low\s+scores?|"
    r"do\s+not\s+recommend|significant\s+concerns?|mixed\s+recommendations?|"
    r"close\s+scores?|important\s+concerns?",
    re.IGNORECASE
)

PUNCTUATION_RE = re.compile(r"[-\u2014\u2013,;/\"\']+")

PROPOSAL_ID_RE = re.compile(r"(PPB\d{2}-\d{2,3})", re.IGNORECASE)


def detect_vote(text):
    """Return vote category from leading emoji, or None."""
    text = text.strip()
    for emoji, label in EMOJI_VOTE_MAP.items():
        if text.startswith(emoji):
            return label
    return None


def parse_notes(notes_text):
    """
    Parse slide notes in the format:
        🟢 Clear yes — description MSO ABC
        🔴 Clear no — description XYZ
        🟡 Needs discussion — description JKL MNO

    Initials are ALL-CAPS tokens at the end of each emoji-prefixed line.
    Returns a list of (initials, vote) tuples.
    """
    results = []
    for line in notes_text.splitlines():
        line = line.strip()
        if not line:
            continue
        vote = detect_vote(line)
        if vote is None:
            continue
        for emoji in EMOJI_VOTE_MAP:
            line = line.replace(emoji, "")
        line = BOILERPLATE_RE.sub(" ", line)
        line = PUNCTUATION_RE.sub(" ", line)
        for token in line.split():
            token = token.strip()
            if re.fullmatch(r"[A-Z]{2,5}", token):
                results.append((token, vote))
    return results


def extract_proposal_id(slide):
    """Return the PPB proposal ID found anywhere on the slide, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                m = PROPOSAL_ID_RE.search(para.text)
                if m:
                    return m.group(1).upper()
    return None


def main():
    # =========================================================================
    # CONFIGURE THIS: paste the path to your OneDrive folder below.
    # This script can live anywhere on your computer -- it does NOT need to
    # be inside OneDrive, so collaborators will never see it.
    #
    # Windows example:
    #   FOLDER_PATH = r"C:\Users\YourName\OneDrive - YourOrg\PPB26 Voting"
    # Mac example:
    #   FOLDER_PATH = "/Users/YourName/Library/CloudStorage/OneDrive-YourOrg/PPB26 Voting"
    # =========================================================================
    FOLDER_PATH = r"/Users/laurasedra/Library/CloudStorage/OneDrive-SharedLibraries-SUNYOldWestbury/Business - Documents"
    # =========================================================================

    if "PASTE YOUR" in FOLDER_PATH:
        sys.exit(
            "ERROR: Please open extract_votes.py and set FOLDER_PATH to your "
            "OneDrive folder path before running."
        )
    if not os.path.isdir(FOLDER_PATH):
        sys.exit(
            "ERROR: Folder not found:\n  {}\n"
            "Check the path is correct and that OneDrive is synced.".format(FOLDER_PATH)
        )

    pptx_files = glob.glob(os.path.join(FOLDER_PATH, "**", "PPB26_Committee_Review.pptx"), recursive=True)
    xlsx_files = glob.glob(os.path.join(FOLDER_PATH, "**", "PPB26_Scores & Rankings.xlsx"), recursive=True)

    if not pptx_files:
        sys.exit("ERROR: No .pptx file found in:\n  {}".format(FOLDER_PATH))
    if not xlsx_files:
        sys.exit("ERROR: No .xlsx file found in:\n  {}".format(FOLDER_PATH))
    if len(pptx_files) > 1:
        sys.exit(
            "ERROR: Multiple .pptx files found: {}\n"
            "Please keep only one .pptx in the folder.".format(pptx_files)
        )
    if len(xlsx_files) > 1:
        sys.exit(
            "ERROR: Multiple .xlsx files found: {}\n"
            "Please keep only one .xlsx in the folder.".format(xlsx_files)
        )

    pptx_path = pptx_files[0]
    xlsx_path = xlsx_files[0]
    print("PowerPoint : {}".format(os.path.basename(pptx_path)))
    print("Excel      : {}".format(os.path.basename(xlsx_path)))

    # Step 1: Parse PowerPoint notes
    prs = Presentation(pptx_path)
    votes_by_proposal = {}
    unmatched_slides = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        proposal_id = extract_proposal_id(slide)

        notes_text = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text if notes_frame else ""

        parsed = parse_notes(notes_text)
        if not parsed:
            continue

        if proposal_id is None:
            unmatched_slides.append((slide_idx, parsed))
            print("  WARNING  Slide {}: no proposal ID found -- votes: {}".format(slide_idx, parsed))
            continue

        pid = proposal_id.upper()
        if pid not in votes_by_proposal:
            votes_by_proposal[pid] = {"YES": [], "NO": [], "NEEDS DISCUSSION": []}
        for initials, vote in parsed:
            votes_by_proposal[pid][vote].append(initials)

    print("\nProposals with votes found: {}".format(list(votes_by_proposal.keys())))

    # Step 2: Find the right sheet in Excel
    wb = load_workbook(xlsx_path)

    target_sheet = None
    for candidate in ("Rank and Average", "Proposal_List", "Summary", "Summ"):
        if candidate in wb.sheetnames:
            target_sheet = wb[candidate]
            break
    if target_sheet is None:
        for name in wb.sheetnames:
            if name not in ("Instructions", "Setup", "Import_Responses", "Scores"):
                target_sheet = wb[name]
                break
    if target_sheet is None:
        target_sheet = wb.active

    print("\nWriting to sheet: '{}'".format(target_sheet.title))

    # Find header row and column positions
    header_row_idx = None
    col_proposal = col_yes = col_no = col_nd = None

    for row in target_sheet.iter_rows():
        for cell in row:
            val = str(cell.value or "").strip().upper()
            if val == "PROPOSAL":
                col_proposal = cell.column
                header_row_idx = cell.row
            elif val == "YES":
                col_yes = cell.column
            elif val == "NO":
                col_no = cell.column
            elif "NEEDS" in val and "DISCUSSION" in val:
                col_nd = cell.column
        if header_row_idx:
            break

    if header_row_idx is None:
        sys.exit(
            "ERROR: Could not find a header row with a 'Proposal' column.\n"
            "Make sure the sheet has headers: Proposal, YES, NO, NEEDS DISCUSSION"
        )

    missing = [
        label for label, col in
        [("YES", col_yes), ("NO", col_no), ("NEEDS DISCUSSION", col_nd)]
        if col is None
    ]
    if missing:
        print("  WARNING: Header columns not found for: {} -- those votes will be skipped.".format(missing))

    # Step 3: Write initials into the correct cells
    def write_initials(row_num, col, vote_list):
        if col and vote_list:
            cell = target_sheet.cell(row=row_num, column=col)
            existing_str = str(cell.value or "").strip()
            # Parse already-recorded initials so we never duplicate them
            already_recorded = set(
                t.strip().upper()
                for t in existing_str.split(",")
                if t.strip()
            )
            new_only = [i for i in vote_list if i.upper() not in already_recorded]
            if not new_only:
                return
            all_initials = sorted(already_recorded | set(i.upper() for i in new_only))
            cell.value = ", ".join(all_initials)

    updated = 0
    for row in target_sheet.iter_rows(min_row=header_row_idx + 1):
        if col_proposal is None:
            break
        proposal_cell = target_sheet.cell(row=row[0].row, column=col_proposal)
        m = PROPOSAL_ID_RE.search(str(proposal_cell.value or ""))
        if not m:
            continue
        pid = m.group(1).upper()
        if pid not in votes_by_proposal:
            continue

        data = votes_by_proposal[pid]
        row_num = row[0].row

        write_initials(row_num, col_yes, data["YES"])
        write_initials(row_num, col_no, data["NO"])
        write_initials(row_num, col_nd, data["NEEDS DISCUSSION"])

        print("  OK  {}: YES={} | NO={} | ND={}".format(
            pid, data["YES"], data["NO"], data["NEEDS DISCUSSION"]
        ))
        updated += 1

    # Step 4: Write tallies into the existing Voting Tally sheet
    tally_sheet_name = "Voting Tally"
    if tally_sheet_name not in wb.sheetnames:
        sys.exit(
            "ERROR: Could not find a sheet named '{}'.\n"
            "Make sure it exists in the Excel file.".format(tally_sheet_name)
        )
    tally = wb[tally_sheet_name]

    # Find header row and YES / NO / NEEDS DISCUSSION columns in tally sheet
    tally_header_row = None
    tally_col_proposal = tally_col_yes = tally_col_no = tally_col_nd = None

    for row in tally.iter_rows():
        for cell in row:
            val = str(cell.value or "").strip().upper()
            if val == "PROPOSAL":
                tally_col_proposal = cell.column
                tally_header_row = cell.row
            elif val == "YES":
                tally_col_yes = cell.column
            elif val == "NO":
                tally_col_no = cell.column
            elif "NEEDS" in val and "DISCUSSION" in val:
                tally_col_nd = cell.column
        if tally_header_row:
            break

    if tally_header_row is None:
        sys.exit(
            "ERROR: Could not find headers in the Voting Tally sheet.\n"
            "Make sure it has columns: Proposal, YES, NO, NEEDS DISCUSSION"
        )

    tally_updated = 0
    for row in tally.iter_rows(min_row=tally_header_row + 1):
        if tally_col_proposal is None:
            break
        proposal_cell = tally.cell(row=row[0].row, column=tally_col_proposal)
        m = PROPOSAL_ID_RE.search(str(proposal_cell.value or ""))
        if not m:
            continue
        pid = m.group(1).upper()
        if pid not in votes_by_proposal:
            continue
        data = votes_by_proposal[pid]
        row_num = row[0].row
        if tally_col_yes:
            tally.cell(row=row_num, column=tally_col_yes, value=len(data["YES"]))
        if tally_col_no:
            tally.cell(row=row_num, column=tally_col_no, value=len(data["NO"]))
        if tally_col_nd:
            tally.cell(row=row_num, column=tally_col_nd, value=len(data["NEEDS DISCUSSION"]))
        tally_updated += 1

    print("\nVoting Tally sheet updated ({} proposals).".format(tally_updated))

    # Save
    out_name = xlsx_path
    wb.save(out_name)
    print("Done. {} proposal(s) updated.".format(updated))
    print("Saved -> {}".format(os.path.basename(out_name)))

    if unmatched_slides:
        print(
            "\nWARNING: {} slide(s) had votes but no detectable proposal ID.\n"
            "Make sure each slide contains the proposal ID (e.g. PPB26-04) "
            "somewhere in the slide text.".format(len(unmatched_slides))
        )


if __name__ == "__main__":
    main()
